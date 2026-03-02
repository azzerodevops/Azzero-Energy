"""Generate PowerPoint (PPTX) presentations using python-pptx.

Creates a branded presentation with:
1. Title slide
2. Site overview
3. KPI summary
4. Technology results
5. Charts (embedded PNGs from matplotlib)
6. Conclusions
"""

from __future__ import annotations

import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from .models import ReportData
from .charts import generate_all_charts

# Brand colours
PRIMARY = RGBColor(0x00, 0x97, 0xD7)
SECONDARY = RGBColor(0x00, 0xB8, 0x94)
DARK_BG = RGBColor(0x12, 0x18, 0x27)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)


def _set_slide_bg(slide, color: RGBColor = DARK_BG) -> None:
    """Set the background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def _add_title_slide(prs: Presentation, data: ReportData) -> None:
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide)

    # Title
    _add_text_box(slide, Inches(1), Inches(2), Inches(8), Inches(1),
                  "AzzeroCO2 Energy", font_size=32, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    _add_text_box(slide, Inches(1), Inches(3), Inches(8), Inches(0.5),
                  "Report Analisi Energetica", font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Analysis info
    info = f"{data.analysis.name} — {data.analysis.site.name}"
    _add_text_box(slide, Inches(1), Inches(4), Inches(8), Inches(0.5),
                  info, font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Date
    _add_text_box(slide, Inches(1), Inches(4.5), Inches(8), Inches(0.5),
                  data.generated_at[:10], font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)


def _add_site_slide(prs: Presentation, data: ReportData) -> None:
    """Slide 2: Site overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                  "Panoramica Sito", font_size=24, color=PRIMARY, bold=True)

    info_lines = [
        f"Sito: {data.analysis.site.name}",
        f"Città: {data.analysis.site.city or '—'}",
        f"Anno di analisi: {data.analysis.year}",
        f"Scenario: {data.scenario.name}",
        f"Obiettivo: {'Minimizzazione costi' if data.scenario.objective == 'cost' else 'Decarbonizzazione'}",
    ]

    y = Inches(1.2)
    for line in info_lines:
        _add_text_box(slide, Inches(1), y, Inches(8), Inches(0.4), line, font_size=14, color=WHITE)
        y += Inches(0.5)

    # Demands summary
    _add_text_box(slide, Inches(0.5), y + Inches(0.3), Inches(9), Inches(0.5),
                  "Consumi Energetici", font_size=18, color=SECONDARY, bold=True)

    y += Inches(0.9)
    for demand in data.demands:
        line = f"• {demand.end_use_label}: {demand.annual_consumption_mwh:,.0f} MWh"
        _add_text_box(slide, Inches(1), y, Inches(8), Inches(0.35), line, font_size=12, color=WHITE)
        y += Inches(0.35)


def _add_kpi_slide(prs: Presentation, data: ReportData) -> None:
    """Slide 3: KPI summary with big numbers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                  "Risultati Chiave", font_size=24, color=PRIMARY, bold=True)

    kpis = [
        ("CAPEX Totale", f"€ {data.scenario.total_capex:,.0f}"),
        ("Risparmio Annuo", f"€ {data.scenario.total_savings_annual:,.0f}"),
        ("Payback", f"{data.scenario.payback_years:.1f} anni" if data.scenario.payback_years else "—"),
        ("Riduzione CO₂", f"{data.scenario.co2_reduction_percent * 100:.1f}%" if data.scenario.co2_reduction_percent else "—"),
    ]

    # 2x2 grid of KPI boxes
    positions = [
        (Inches(0.5), Inches(1.5)),
        (Inches(5), Inches(1.5)),
        (Inches(0.5), Inches(3.5)),
        (Inches(5), Inches(3.5)),
    ]

    for (left, top), (label, value) in zip(positions, kpis):
        # Box background
        shape = slide.shapes.add_shape(
            1, left, top, Inches(4.3), Inches(1.5)  # 1 = rectangle
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()  # no border

        # Label
        _add_text_box(slide, left + Inches(0.3), top + Inches(0.2), Inches(3.7), Inches(0.4),
                      label, font_size=12, color=MUTED)
        # Value
        _add_text_box(slide, left + Inches(0.3), top + Inches(0.6), Inches(3.7), Inches(0.7),
                      value, font_size=28, color=WHITE, bold=True)


def _add_tech_slide(prs: Presentation, data: ReportData) -> None:
    """Slide 4: Technology results table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                  "Tecnologie Ottimali", font_size=24, color=PRIMARY, bold=True)

    y = Inches(1.3)
    for tech in data.scenario.tech_results:
        line = (
            f"• {tech.technology_name} ({tech.category}) — "
            f"{tech.optimal_capacity_kw:,.0f} kW | "
            f"CAPEX €{tech.capex:,.0f} | "
            f"Risparmio €{tech.annual_savings:,.0f}/anno"
        )
        _add_text_box(slide, Inches(0.5), y, Inches(9), Inches(0.4), line, font_size=11, color=WHITE)
        y += Inches(0.45)


def _add_chart_slide(prs: Presentation, title: str, chart_path: Path) -> None:
    """Add a slide with an embedded chart image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                  title, font_size=24, color=PRIMARY, bold=True)

    if chart_path.exists():
        slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1.2), width=Inches(9))


def _add_conclusions_slide(prs: Presentation, data: ReportData) -> None:
    """Final slide with conclusions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_text_box(slide, Inches(1), Inches(2), Inches(8), Inches(0.7),
                  "Conclusioni", font_size=28, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

    total_demand = sum(d.annual_consumption_mwh for d in data.demands)
    total_production = sum(t.annual_production_mwh for t in data.scenario.tech_results)

    conclusions = [
        f"Investimento totale: € {data.scenario.total_capex:,.0f}",
        f"Risparmio annuo: € {data.scenario.total_savings_annual:,.0f}",
        f"Copertura del fabbisogno: {(total_production / total_demand * 100):.0f}%" if total_demand > 0 else "—",
    ]
    if data.scenario.payback_years:
        conclusions.append(f"Recupero investimento in {data.scenario.payback_years:.1f} anni")
    if data.scenario.co2_reduction_percent:
        conclusions.append(f"Riduzione emissioni CO₂: {data.scenario.co2_reduction_percent * 100:.1f}%")

    y = Inches(3)
    for c in conclusions:
        _add_text_box(slide, Inches(1), y, Inches(8), Inches(0.4),
                      f"• {c}", font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
        y += Inches(0.45)

    _add_text_box(slide, Inches(1), Inches(6.5), Inches(8), Inches(0.3),
                  "AzzeroCO2 Energy — Il clima nelle nostre mani",
                  font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER)


def generate_pptx(data: ReportData) -> Path:
    """Generate a PowerPoint presentation from report data.

    Returns path to the generated .pptx file.
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Generate charts
    charts = generate_all_charts(data)

    # Build slides
    _add_title_slide(prs, data)
    _add_site_slide(prs, data)
    _add_kpi_slide(prs, data)
    _add_tech_slide(prs, data)

    # Chart slides
    chart_titles = {
        "energy_mix": "Mix Energetico",
        "capex_breakdown": "Ripartizione CAPEX",
        "capex_vs_savings": "CAPEX vs Risparmio",
        "cashflow": "Flusso di Cassa Cumulato",
    }
    for key, title in chart_titles.items():
        if key in charts:
            _add_chart_slide(prs, title, charts[key])

    _add_conclusions_slide(prs, data)

    # Save
    output_dir = Path(tempfile.mkdtemp(prefix="azzeroco2_report_"))
    filename = f"report_{data.analysis.name.replace(' ', '_')}_{data.generated_at[:10]}.pptx"
    output_path = output_dir / filename

    prs.save(str(output_path))
    return output_path
